#!/usr/bin/env python
# coding: utf-8

# ### Hier Dokumentiert, was in den einzelnen Schritten gemacht wird. Einzelne Abschnitte werden mit Markdown kommentiert, damit es übersichtlicher wird. 
# Die restliche Dokumentation wird über Kommentare in den Codeblöcken gemacht. 
# Die angegbenen Links sind Quellen aus denen ich Codes teilweise übernommen oder die als inspiration gedient haben.

# In[1]:


#https://python-pptx.readthedocs.io/en/latest/user/quickstart.html
from pptx import Presentation
from pptx.util import Inches
from pptx.util import Cm
#Maße PPTX Folien WIDTH=960 x HEIGHT=720 Pixel. in Inches = WIDTH=10 x HEIGHT=7.5 WIDTH= 8,13cm x HEIGHT=6,10cm

img_path = "report_head_ppx.png"

prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "TED Talk Explorer - Report"
subtitle.text = "In diesem Report werden die Ergebnisse der Datenauswertung detaillierter dargestellt."

